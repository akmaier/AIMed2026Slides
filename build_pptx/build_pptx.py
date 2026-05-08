#!/usr/bin/env python3
"""Convert the rendered PDF deck to a PowerPoint, embedding animated GIFs
on the three "What is agentic AI?" attempt slides.

For every PDF page we add a full-bleed picture of the rasterized page.
On pages 25/26/27 we additionally drop the source GIF on top of the
static PNG so PowerPoint shows the animation in slideshow mode.

The GIF placement matches the static image's bounding box on the
rasterized page exactly:

  1. Find the bbox of non-white content inside the source PNG file
     (the 1200x1200 image has whitespace around the matplotlib figure).
  2. Find the bbox of non-white content inside the rasterized page's
     body region (skipping header & footer).
  3. The ratio between the two gives the on-slide scale factor; from
     that we recover the full 1200x1200 placement rectangle and put the
     GIF there. Since GIF and PNG are the same size and layout, the
     GIF aligns pixel-for-pixel with the static image.
"""

from __future__ import annotations

import subprocess
from pathlib import Path

import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Emu

ROOT = Path(__file__).resolve().parent.parent
PDF = ROOT / "output_tex" / "main.pdf"
FIG_DIR = ROOT / "output_tex" / "figures"
PAGES_DIR = ROOT / "build_pptx" / "pages"
# Write the regenerated deck *inside* build_pptx/ so we never clobber
# the manually-edited committed copy at the repo root.
OUT = ROOT / "build_pptx" / "AIMed2026_What_next_in_medical_AI.pptx"

# Slide-page -> (gif filename, source png filename)
TARGETS = {
    25: ("img_08_02.gif", "img_08_02.png"),
    26: ("img_08_03.gif", "img_08_03.png"),
    27: ("img_08_04.gif", "img_08_04.png"),
}

# Body region of an FAU-themed slide (fractions of page height).
# Tightened to skip the title bar / logos (top) and footer (bottom);
# determined empirically from probe_layout.py on a target page.
BODY_TOP_FRAC = 0.22
BODY_BOT_FRAC = 0.95


def content_bbox(arr: np.ndarray, white_threshold: int = 240) -> tuple[int, int, int, int] | None:
    """(xmin, ymin, xmax, ymax) of pixels darker than `white_threshold`."""
    if arr.ndim == 3:
        # Handle RGBA: treat fully transparent as white
        if arr.shape[2] == 4:
            rgb = arr[..., :3].astype(np.float32)
            alpha = arr[..., 3:4].astype(np.float32) / 255.0
            arr = (rgb * alpha + 255 * (1 - alpha)).astype(np.uint8)
        gray = arr.mean(axis=2)
    else:
        gray = arr
    mask = gray < white_threshold
    rows = np.any(mask, axis=1)
    cols = np.any(mask, axis=0)
    if not rows.any() or not cols.any():
        return None
    rmin, rmax = np.where(rows)[0][[0, -1]]
    cmin, cmax = np.where(cols)[0][[0, -1]]
    return int(cmin), int(rmin), int(cmax) + 1, int(rmax) + 1


def gif_placement(page_png: Path, source_png: Path,
                  page_w: int, page_h: int) -> tuple[float, float, float, float]:
    """Return (x_frac, y_frac, w_frac, h_frac) for the GIF on the slide,
    in fractions of slide width/height."""
    # Source PNG content bbox (in PNG pixels)
    src = np.array(Image.open(source_png))
    src_h, src_w = src.shape[:2]
    sb = content_bbox(src)
    if sb is None:
        raise RuntimeError(f"No content found in {source_png}")
    sx0, sy0, sx1, sy1 = sb
    src_cw, src_ch = sx1 - sx0, sy1 - sy0

    # Rasterized page: content bbox restricted to the body region
    page = np.array(Image.open(page_png))
    pH, pW = page.shape[:2]
    body_top = int(pH * BODY_TOP_FRAC)
    body_bot = int(pH * BODY_BOT_FRAC)
    pb = content_bbox(page[body_top:body_bot, :])
    if pb is None:
        raise RuntimeError(f"No body content found in {page_png}")
    px0, py0_rel, px1, py1_rel = pb
    py0 = py0_rel + body_top
    py1 = py1_rel + body_top
    page_cw, page_ch = px1 - px0, py1 - py0

    # Scale: rasterized content size / source content size
    scale_w = page_cw / src_cw
    scale_h = page_ch / src_ch

    # Full PNG placement on the rasterized page (in page pixels)
    full_x = px0 - sx0 * scale_w
    full_y = py0 - sy0 * scale_h
    full_w = src_w * scale_w
    full_h = src_h * scale_h

    return full_x / pW, full_y / pH, full_w / pW, full_h / pH


def main() -> None:
    page_files = sorted(PAGES_DIR.glob("page-*.png"))
    if not page_files:
        raise SystemExit(f"No rasterized pages in {PAGES_DIR}; "
                         "run pdftoppm first.")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for i, page_png in enumerate(page_files, start=1):
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(page_png), 0, 0,
            width=prs.slide_width, height=prs.slide_height,
        )

        if i in TARGETS:
            gif_name, png_name = TARGETS[i]
            with Image.open(page_png) as im:
                pw, ph = im.size
            xf, yf, wf, hf = gif_placement(
                page_png, FIG_DIR / png_name, pw, ph,
            )
            x = int(prs.slide_width * xf)
            y = int(prs.slide_height * yf)
            w = int(prs.slide_width * wf)
            h = int(prs.slide_height * hf)
            slide.shapes.add_picture(
                str(FIG_DIR / gif_name), x, y, width=w, height=h,
            )
            print(f"  page {i}: overlaid {gif_name} "
                  f"x={xf:.4f} y={yf:.4f} w={wf:.4f} h={hf:.4f}")

    prs.save(str(OUT))
    print(f"\nWrote {OUT} ({OUT.stat().st_size:,} bytes, "
          f"{len(page_files)} slides)")


if __name__ == "__main__":
    main()
